import sys
import subprocess

def install(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])

try:
    import docx
except ImportError:
    install('python-docx')
    import docx

try:
    import pptx
except ImportError:
    install('python-pptx')
    import pptx

from docx import Document
from pptx import Presentation

# Ensure paths are correct
report_path = 'c:\\Users\\acer\\OneDrive\\Desktop\\Agri-Mitra-Hub-Source\\Agri-Mitra-Hub_MidTerm_Report.docx'
ppt_path = 'c:\\Users\\acer\\OneDrive\\Desktop\\Agri-Mitra-Hub-Source\\Agri-Mitra-Hub_MidTerm_Presentation.pptx'

print("Generating Document...", flush=True)
# Create Report
doc = Document()
doc.add_heading('Agri-Mitra-Hub: Mid-Term Project Report', 0)
doc.add_heading('1. Introduction', level=1)
doc.add_paragraph('Agri-Mitra-Hub is an integrated platform empowering Indian farmers by providing real-time market prices, AI-driven crop recommendations, and smart irrigation tracking.')
doc.add_heading('2. Project Objectives', level=1)
doc.add_paragraph('- Provide Live Mandi Prices for all crops across states.\n- Yield Estimator to calculate expected profits.\n- Farm Profiles for managing soil history.\n- Transport Finder for estimating transport cost to nearby mandis.\n- Mandi Locator via GPS mapping.\n- Weather Alerts for rain, frost, pest, and heatwave alerts.\n- Modern Techniques Guide with training videos.')
doc.add_heading('3. Current Progress', level=1)
doc.add_paragraph('The frontend architecture is completed using React, TailwindCSS, and Vite. Core routing and the dashboard UI structure mapping out the Quick Access modules are fully designed and functional. Backend integration is underway with Express.')
doc.add_heading('4. Future Scope', level=1)
doc.add_paragraph('Future milestones involve deploying the full AI engine (Yield Estimator), hooking up real-time hardware or API sensors for weather and soil status, and integrating secure payment and SMS alert systems for offline farmers.')
doc.save(report_path)

print("Generating Presentation...", flush=True)
# Create Presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Agri-Mitra-Hub"
subtitle.text = "Mid-Term Project Presentation\nEmpowering Indian Farmers with AI & Real-time Data"

bullet_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(bullet_slide_layout)
shapes2 = slide2.shapes
title_shape2 = shapes2.title
body_shape2 = shapes2.placeholders[1]
title_shape2.text = "Core Objectives"
tf2 = body_shape2.text_frame
tf2.text = "Provide real-time Mandi market prices"
p = tf2.add_paragraph()
p.text = "Yield Estimator for profit calculation"
p = tf2.add_paragraph()
p.text = "Transport Finder & GPS Mandi Locator"
p = tf2.add_paragraph()
p.text = "Weather & Pest control alerts"

slide3 = prs.slides.add_slide(bullet_slide_layout)
shapes3 = slide3.shapes
title_shape3 = shapes3.title
body_shape3 = shapes3.placeholders[1]
title_shape3.text = "Implementation Status"
tf3 = body_shape3.text_frame
tf3.text = "Modern dashboard built with Vite + React"
p = tf3.add_paragraph()
p.text = "Backend architecture defined using Express"
p = tf3.add_paragraph()
p.text = "Successfully embedded training and modern techniques guide"
prs.save(ppt_path)

print("SUCCESSFULLY GENERATED BOTH FILES")
